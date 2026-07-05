from __future__ import annotations

import re
from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
PAPER = ROOT / "paper"
TEMPLATE = PAPER / "SWAIB_TransG_slides.pptx"
SPEECH = PAPER / "cmtflow_20min_speech.md"
OUT = PAPER / "CMTFlow_20min_editable_draft.pptx"

FONT = "微软雅黑"
FONT_EN = "Georgia"

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(90, 90, 90)
LIGHT_GRAY = RGBColor(244, 246, 249)
MID_GRAY = RGBColor(210, 218, 230)
RED = RGBColor(192, 0, 0)
BLUE = RGBColor(0, 112, 192)
PURPLE = RGBColor(112, 48, 160)
GREEN = RGBColor(0, 176, 80)
NAVY = RGBColor(31, 73, 125)
SOFT_RED = RGBColor(252, 235, 235)
SOFT_BLUE = RGBColor(232, 242, 252)
SOFT_PURPLE = RGBColor(241, 235, 248)
SOFT_GREEN = RGBColor(235, 247, 239)


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def parse_speech(path: Path) -> dict[int, str]:
    text = path.read_text(encoding="utf-8")
    sections = re.split(r"(?m)^## Slide (\d+): [^\n]+\n", text)
    notes: dict[int, str] = {}
    for i in range(1, len(sections), 2):
        slide_no = int(sections[i])
        body = sections[i + 1].strip()
        notes[slide_no] = body
    return notes


def set_run_font(run, size=18, color=BLACK, bold=False, italic=False, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def set_cell_text(cell, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        set_run_font(run, size=size, color=color, bold=bold)


def add_textbox(
    slide,
    text,
    left,
    top,
    width,
    height,
    size=18,
    color=BLACK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font=FONT,
    margin=0.05,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        set_run_font(run, size=size, color=color, bold=bold, italic=italic, name=font)
    return box


def add_bullets(slide, bullets, left, top, width, height, size=16, color=BLACK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        for run in p.runs:
            set_run_font(run, size=size, color=color)
    return box


def add_title(slide, title, section=None):
    add_textbox(slide, title, 0.78, 0.34, 11.65, 0.55, size=26, color=BLACK, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(1.05), Inches(2.1), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    if section:
        add_textbox(slide, section, 10.3, 0.28, 2.25, 0.38, size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_footer(slide, n, label="CMTFlow"):
    add_textbox(slide, label, 0.78, 7.05, 2.0, 0.22, size=8, color=GRAY)
    add_textbox(slide, f"{n:02d}", 12.15, 7.05, 0.45, 0.22, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def add_card(slide, title, body, left, top, width, height, fill, accent=RED, title_size=15, body_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.1)
    add_textbox(slide, title, left + 0.14, top + 0.12, width - 0.28, 0.32, size=title_size, color=accent, bold=True)
    add_textbox(slide, body, left + 0.14, top + 0.55, width - 0.28, height - 0.68, size=body_size, color=BLACK)
    return shape


def add_pill(slide, text, left, top, width, color, fill=WHITE, size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    add_textbox(slide, text, left + 0.04, top + 0.04, width - 0.08, 0.22, size=size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return shape


def add_arrow(slide, left, top, width, height, color=GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_picture_fit(slide, path: Path, left, top, width, height, border=True):
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = width / height
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        w = width
        h = width / img_ratio
    else:
        h = height
        w = height * img_ratio
    x = left + (width - w) / 2
    y = top + (height - h) / 2
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if border:
        pic.line.color.rgb = MID_GRAY
        pic.line.width = Pt(0.75)
    return pic


def add_notes(slide, notes: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = notes


def new_slide(prs, slide_no, title, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, title, section=section)
    add_footer(slide, slide_no)
    return slide


def cover(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.18)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = RED
    slide.shapes[-1].line.fill.background()
    add_textbox(slide, "CMTFlow", 0.9, 1.25, 4.0, 0.72, size=38, color=RED, bold=True, font=FONT_EN)
    add_textbox(slide, "控制器引导的分层投资组合管理框架", 0.9, 2.05, 10.7, 0.75, size=28, color=BLACK, bold=True)
    add_textbox(slide, "Hierarchical Portfolio Management with Controller-Guided Base Revision and Daily Refinement", 0.92, 2.88, 10.7, 0.42, size=14, color=GRAY, font=FONT_EN)
    add_card(slide, "报告主线", "何时修正当前持仓  |  换成什么基准组合  |  如何每日局部微调", 0.95, 4.05, 9.6, 0.92, SOFT_BLUE, BLUE, title_size=13, body_size=15)
    add_textbox(slide, "作者 / 单位 / 日期", 0.95, 5.75, 6.5, 0.4, size=14, color=GRAY)
    add_footer(slide, 1)
    add_notes(slide, notes.get(1, ""))


def build():
    notes = parse_speech(SPEECH)
    prs = Presentation(TEMPLATE)
    delete_all_slides(prs)
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    cover(prs, notes)

    s = new_slide(prs, 2, "研究背景：投资组合管理的基本目标", "Background")
    add_bullets(s, ["多资产资本配置：在不确定市场中分配资金", "目标不是单纯收益最大化，而是收益、风险、成本的长期权衡", "强化学习适合序列决策，但金融市场具有持续非平稳性"], 0.98, 1.55, 5.0, 2.35, size=17)
    add_card(s, "收益", "Total Return\nAnnual Return", 6.35, 1.65, 1.75, 1.2, SOFT_RED, RED)
    add_card(s, "风险", "Volatility\nMax Drawdown", 8.45, 1.65, 1.75, 1.2, SOFT_BLUE, BLUE)
    add_card(s, "成本", "Turnover\nTransaction Cost", 10.55, 1.65, 1.75, 1.2, SOFT_PURPLE, PURPLE)
    add_arrow(s, 6.0, 3.45, 6.0, 0.42, MID_GRAY)
    add_textbox(s, "长期投资路径质量", 7.55, 4.02, 3.35, 0.45, size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(2, ""))

    s = new_slide(prs, 3, "研究背景：组合决策不是单一日频动作", "Background")
    steps = [("中期持仓", "形成 base portfolio\n稳定持仓方向", RED), ("每日修正", "根据新信息微调\n控制短期偏离", BLUE), ("异常退出", "状态恶化时修正\n避免陈旧持仓", PURPLE)]
    for i, (t, b, c) in enumerate(steps):
        add_card(s, t, b, 1.1 + i * 3.85, 2.0, 2.75, 1.55, [SOFT_RED, SOFT_BLUE, SOFT_PURPLE][i], c, title_size=17, body_size=13)
        if i < 2:
            add_arrow(s, 3.95 + i * 3.85, 2.48, 0.62, 0.38, MID_GRAY)
    add_textbox(s, "固定周期再平衡：稳定但反应慢", 1.18, 4.55, 4.8, 0.42, size=17, color=GRAY, bold=True)
    add_textbox(s, "纯日频调仓：灵活但易受噪声驱动", 6.4, 4.55, 5.2, 0.42, size=17, color=GRAY, bold=True)
    add_notes(s, notes.get(3, ""))

    s = new_slide(prs, 4, "问题定义：带漂移和交易成本的动态组合", "Problem Definition")
    flow = [("w_{t-1}", "昨日组合"), ("Drift", "价格变化导致权重漂移"), ("Rebalance", "主动调整产生换手"), ("w_t", "最终执行权重"), ("Reward", "log return - cost")]
    for i, (t, b) in enumerate(flow):
        add_card(s, t, b, 0.85 + i * 2.45, 2.0, 1.75, 1.18, WHITE, [RED, BLUE, PURPLE, GREEN, NAVY][i], title_size=16, body_size=10)
        if i < len(flow) - 1:
            add_arrow(s, 2.68 + i * 2.45, 2.38, 0.44, 0.28, MID_GRAY)
    add_bullets(s, ["组合会自然漂移：即使不交易，资产价格变化也会改变真实暴露", "重新配置会带来交易成本，需要和收益共同优化", "评价指标同时关注 TR、AR、Sharpe、MDD、CR"], 1.0, 4.2, 10.8, 1.45, size=16)
    add_notes(s, notes.get(4, ""))

    s = new_slide(prs, 5, "问题定义：本文关注的三个核心决策", "Problem Definition")
    cards = [
        ("When", "当前 base portfolio 是否已经失效？\n是否应该退出当前持仓段？", RED),
        ("What", "如果切换，下一段应该持有\n什么样的基准组合？", BLUE),
        ("How", "在选定基准组合内部，\n如何进行每日权重微调？", PURPLE),
    ]
    for i, (t, b, c) in enumerate(cards):
        add_card(s, t, b, 1.0 + i * 4.05, 1.8, 3.1, 2.0, [SOFT_RED, SOFT_BLUE, SOFT_PURPLE][i], c, title_size=24, body_size=14)
        add_textbox(s, ["Controller", "Outer Actor", "Inner Actor"][i], 1.08 + i * 4.05, 4.05, 2.9, 0.35, size=15, color=c, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, "核心思想：把混杂的每日动作拆成不同时间尺度上的协同决策", 1.2, 5.35, 10.8, 0.5, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(5, ""))

    s = new_slide(prs, 6, "相关工作：从静态优化到深度强化学习", "Related Work")
    cols = [
        ("Classical Optimization", "均值方差、CAPM、规则再平衡\n解释性强，但依赖静态假设或单期目标", RED),
        ("Deep RL Portfolio", "直接优化长期收益\n常把中期配置和每日执行压缩成一个动作", BLUE),
        ("Adaptive / HRL", "开始拆分决策角色\n但缺少可持有、漂移、替换的 base memory", PURPLE),
    ]
    for i, (t, b, c) in enumerate(cols):
        add_card(s, t, b, 0.9 + i * 4.0, 1.7, 3.25, 2.45, [SOFT_RED, SOFT_BLUE, SOFT_PURPLE][i], c, title_size=15, body_size=13)
    add_textbox(s, "本文切入点：显式学习 base portfolio 的修正时机", 1.0, 5.2, 11.2, 0.5, size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(6, ""))

    s = new_slide(prs, 7, "挑战：直接应用 RL 仍存在非平凡限制", "Challenge")
    challenges = [
        ("1", "日频动作容易被短期噪声驱动"),
        ("2", "固定再平衡无法识别持仓恶化"),
        ("3", "换仓必须同时比较旧组合、候选组合、成本和持仓年龄"),
    ]
    for i, (num, text) in enumerate(challenges):
        add_card(s, num, text, 1.0, 1.65 + i * 1.25, 5.0, 0.9, WHITE, [RED, BLUE, PURPLE][i], title_size=22, body_size=14)
    add_card(s, "Controller", "学习 hold / switch\n把日历规则变成状态依赖事件策略", 7.0, 2.25, 4.8, 2.05, SOFT_BLUE, BLUE, title_size=22, body_size=16)
    add_textbox(s, "?", 6.25, 2.45, 0.55, 0.8, size=42, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(7, ""))

    s = new_slide(prs, 8, "方案总览：CMTFlow 的统一分层结构", "Overview")
    add_picture_fit(s, PAPER / "figures/cmtflow_architecture_vector.png", 0.9, 1.35, 8.25, 5.15)
    add_bullets(s, ["Outer：生成候选 base portfolio", "Controller：判断 hold 或 switch", "Inner：围绕当前 base 做每日微调", "环境反馈收益、成本和训练信号"], 9.45, 1.55, 3.0, 3.6, size=14)
    add_notes(s, notes.get(8, ""))

    s = new_slide(prs, 9, "强化学习建模：状态、动作与奖励", "Formulation")
    nodes = [("State", "近期市场张量\n漂移持仓\n候选/比较特征"), ("Action", "候选 base\nhold/switch\n执行权重"), ("Environment", "价格变化\n交易成本\n组合收益"), ("Reward", "log return\n- transaction cost")]
    for i, (t, b) in enumerate(nodes):
        add_card(s, t, b, 0.95 + i * 3.05, 2.1, 2.2, 1.55, WHITE, [RED, BLUE, PURPLE, GREEN][i], title_size=16, body_size=12)
        if i < 3:
            add_arrow(s, 3.22 + i * 3.05, 2.65, 0.48, 0.32, MID_GRAY)
    add_textbox(s, "目标：优化完整投资路径上的风险收益质量，而不是单日预测准确率", 1.0, 4.65, 11.2, 0.55, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(9, ""))

    s = new_slide(prs, 10, "方法架构：从候选组合到最终执行权重", "Methodology")
    add_picture_fit(s, PAPER / "figures/cmtflow_decision_flow_imagegen.png", 0.9, 1.32, 8.2, 5.1)
    add_card(s, "决策规则", "b_t = w_t^{cand} if switch\nelse \\tilde{b}_t\n\nw_t = Inner(b_t, state_t)", 9.45, 1.55, 3.0, 1.9, SOFT_BLUE, BLUE, title_size=17, body_size=15)
    add_bullets(s, ["Controller 比较旧 base 与候选 base", "Base portfolio 可持有、漂移、替换", "Inner 始终输出最终执行权重"], 9.48, 3.8, 2.95, 1.85, size=13)
    add_notes(s, notes.get(10, ""))

    s = new_slide(prs, 11, "方法细节：Outer Actor 生成中期基准组合", "Methodology")
    add_picture_fit(s, PAPER / "figures/cmtflow_architecture_vector.png", 0.85, 1.35, 5.75, 4.8)
    add_card(s, "Input", "长期市场窗口\n当前持仓漂移状态", 7.0, 1.55, 2.2, 1.2, SOFT_RED, RED)
    add_arrow(s, 9.35, 2.02, 0.55, 0.28, MID_GRAY)
    add_card(s, "Outer", "segment-level\nasset selector", 10.0, 1.55, 2.2, 1.2, SOFT_BLUE, BLUE)
    add_arrow(s, 8.1, 3.35, 0.55, 0.28, MID_GRAY)
    add_card(s, "Output", "稀疏 top-K\ncandidate base portfolio", 8.85, 3.0, 2.9, 1.35, SOFT_PURPLE, PURPLE)
    add_textbox(s, "Outer 回答 What to hold next，而不是 How to trade today", 7.0, 5.35, 5.2, 0.5, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(11, ""))

    s = new_slide(prs, 12, "方法细节：Controller 学习何时替换基准组合", "Methodology")
    add_picture_fit(s, PAPER / "figures/cmtflow_decision_flow_imagegen.png", 0.85, 1.38, 5.95, 4.7)
    add_card(s, "Controller 输入", "近期市场张量\n漂移当前持仓\n候选组合\nholding/action features", 7.1, 1.45, 2.55, 2.2, SOFT_BLUE, BLUE, title_size=15, body_size=12)
    add_card(s, "Controller 输出", "exit probability\nhold / switch", 10.0, 1.45, 2.2, 1.2, SOFT_RED, RED, title_size=15, body_size=13)
    add_card(s, "最终评价", "每日检查\n无最小持仓约束\n30 天最大持仓上限", 9.2, 3.35, 2.75, 1.65, SOFT_PURPLE, PURPLE, title_size=15, body_size=12)
    add_textbox(s, "把固定日历调仓变成状态依赖的事件策略", 7.1, 5.45, 5.1, 0.42, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(12, ""))

    s = new_slide(prs, 13, "方法细节：Inner Actor 做每日局部微调", "Methodology")
    add_picture_fit(s, PAPER / "figures/cmtflow_architecture_vector.png", 0.85, 1.35, 6.2, 4.9)
    add_card(s, "Base portfolio", "中期持仓锚点\n由 Controller 选择", 7.55, 1.75, 2.25, 1.15, SOFT_BLUE, BLUE)
    add_arrow(s, 9.95, 2.12, 0.55, 0.28, MID_GRAY)
    add_card(s, "Inner tilt", "executed weight\n- base weight", 10.65, 1.75, 1.85, 1.15, SOFT_PURPLE, PURPLE)
    add_card(s, "定位", "局部调权模块\n提供日频灵活性\n不是主要换仓控制器", 7.55, 3.45, 4.8, 1.55, SOFT_GREEN, GREEN, title_size=16, body_size=13)
    add_notes(s, notes.get(13, ""))

    s = new_slide(prs, 14, "模型训练：固定 HRL 预训练与 Controller 学习", "Training")
    add_picture_fit(s, PAPER / "figures/cmtflow_training_flow_vector.png", 0.85, 1.35, 8.2, 5.05)
    phases = [("1", "Fixed HRL warmup\n训练 Outer / Inner"), ("2", "Controller learning\n学习每日 hold / switch"), ("3", "Final evaluation\n每日事件触发")]
    for i, (n, text) in enumerate(phases):
        add_card(s, n, text, 9.45, 1.45 + i * 1.55, 2.85, 1.05, [SOFT_RED, SOFT_BLUE, SOFT_GREEN][i], [RED, BLUE, GREEN][i], title_size=18, body_size=11)
    add_notes(s, notes.get(14, ""))

    s = new_slide(prs, 15, "实验设置", "Experiments")
    table = s.shapes.add_table(3, 5, Inches(0.85), Inches(1.45), Inches(7.65), Inches(1.55)).table
    headers = ["Market", "#Stocks", "Train", "Valid", "Test"]
    rows = [
        ["Nasdaq-100", "39", "2000-04-07\n2017-12-29", "2018-01-02\n2020-04-22", "2020-04-23\n2025-10-03"],
        ["CSI-300", "53", "2000-04-07\n2017-12-28", "2018-01-02\n2019-12-31", "2020-01-02\n2025-02-28"],
    ]
    for j, h in enumerate(headers):
        set_cell_text(table.cell(0, j), h, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        table.cell(0, j).fill.solid()
        table.cell(0, j).fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            set_cell_text(table.cell(i, j), val, size=9.5, align=PP_ALIGN.CENTER)
            table.cell(i, j).fill.solid()
            table.cell(i, j).fill.fore_color.rgb = WHITE if i == 1 else LIGHT_GRAY
    add_card(s, "Baselines", "传统策略\n深度 RL baseline\n固定 5/10/20/30/60 天 controller", 9.0, 1.45, 3.25, 1.55, SOFT_BLUE, BLUE)
    add_card(s, "Metrics", "TR / AR / Vol\nSharpe ratio\nMDD / CR", 9.0, 3.35, 3.25, 1.55, SOFT_PURPLE, PURPLE)
    add_textbox(s, "Transaction cost rate = 5e-5  |  Outer window = 60  |  Inner window = 10  |  Controller window = 30", 0.9, 5.65, 11.7, 0.42, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(15, ""))

    s = new_slide(prs, 16, "数值结果：主实验性能对比", "Results")
    add_picture_fit(s, PAPER / "figures/main_equity_curves.png", 0.75, 1.28, 5.95, 4.9)
    add_picture_fit(s, PAPER / "figures/main_metric_bars.png", 6.95, 1.28, 5.55, 4.9)
    add_textbox(s, "Nasdaq: TR 265.53%, MDD 18.62%  |  CSI-300: TR 204.99%, Sharpe 1.14, MDD 22.78%", 0.9, 6.35, 11.6, 0.38, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(16, ""))

    s = new_slide(prs, 17, "数值结果：消融实验与机制验证", "Results")
    add_picture_fit(s, PAPER / "figures/ablation_metric_bars.png", 0.85, 1.35, 7.7, 4.95)
    add_card(s, "Controller 贡献", "Nasdaq:\nOuter-only 220.42% ->\nOuter + Controller 237.50%\nMDD 32.09% -> 21.24%", 9.0, 1.35, 3.3, 1.9, SOFT_BLUE, BLUE, body_size=11)
    add_card(s, "固定窗口对照", "5/10/20/30/60 天切仓\n不能稳定复现 Ours\n优势来自状态依赖时机", 9.0, 3.55, 3.3, 1.55, SOFT_RED, RED, body_size=11)
    add_notes(s, notes.get(17, ""))

    s = new_slide(prs, 18, "案例研究：Controller 与 Inner 的可解释行为", "Case Study")
    add_picture_fit(s, PAPER / "figures/explainability/controller_switch_cases.png", 0.75, 1.28, 4.0, 4.75)
    add_picture_fit(s, PAPER / "figures/explainability/random_switch_comparison.png", 4.95, 1.28, 3.75, 4.75)
    add_picture_fit(s, PAPER / "figures/explainability/inner_actor_base_adjustment.png", 8.9, 1.28, 3.65, 4.75)
    add_textbox(s, "关键风险窗口：switch 改善 20 日反事实收益并降低回撤；随机切换无法替代 learned controller；Inner 提供局部 tilt。", 0.9, 6.28, 11.5, 0.45, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, notes.get(18, ""))

    s = new_slide(prs, 19, "讨论与总结", "Conclusion")
    summaries = [
        ("1", "问题重构", "base revision\nbase construction\ndaily refinement"),
        ("2", "核心机制", "Controller 将固定周期换仓\n变成可学习事件策略"),
        ("3", "实验结论", "更稳健的风险收益权衡\nController 是主要自适应来源"),
    ]
    for i, (n, t, b) in enumerate(summaries):
        add_card(s, f"{n}. {t}", b, 0.95 + i * 4.05, 1.85, 3.1, 2.05, [SOFT_RED, SOFT_BLUE, SOFT_PURPLE][i], [RED, BLUE, PURPLE][i], title_size=18, body_size=14)
    add_card(s, "未来工作", "更多市场与资产类别  |  更复杂交易约束  |  更稳健风险控制", 1.0, 5.0, 11.25, 0.82, SOFT_GREEN, GREEN, title_size=15, body_size=13)
    add_notes(s, notes.get(19, ""))

    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_textbox(s, "谢谢，欢迎提问", 0.95, 2.25, 11.3, 0.85, size=36, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, "Learning when to revise is as important as learning what to hold.", 1.25, 3.35, 10.8, 0.45, size=18, color=NAVY, italic=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_textbox(s, "CMTFlow", 5.25, 4.65, 2.8, 0.5, size=24, color=BLACK, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_footer(s, 20)
    add_notes(s, notes.get(20, ""))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
