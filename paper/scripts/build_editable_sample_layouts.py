from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_cmtflow_editable_ppt import (
    BLACK,
    BLUE,
    FONT_EN,
    GRAY,
    GREEN,
    LIGHT_GRAY,
    MID_GRAY,
    NAVY,
    PURPLE,
    RED,
    SOFT_BLUE,
    SOFT_GREEN,
    SOFT_PURPLE,
    SOFT_RED,
    TEMPLATE,
    WHITE,
    add_arrow,
    add_card,
    add_textbox,
    delete_all_slides,
    set_run_font,
)


PAPER = Path(__file__).resolve().parents[1]
OUT = PAPER / "cmtflow_codex_ppt/sample_layouts_editable.pptx"


def placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(idx)


def set_placeholder_text(ph, text, size=26, color=BLACK, bold=True, align=PP_ALIGN.LEFT):
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        set_run_font(run, size=size, color=color, bold=bold)


def footer(slide, no):
    for idx, text, align in [
        (10, "", PP_ALIGN.LEFT),
        (11, "CMTFlow editable sample", PP_ALIGN.CENTER),
        (12, f"{no:02d}", PP_ALIGN.RIGHT),
    ]:
        try:
            ph = placeholder(slide, idx)
        except KeyError:
            continue
        set_placeholder_text(ph, text, size=9, color=GRAY, bold=False, align=align)


def title_slide(prs, no, title, layout=5):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    set_placeholder_text(placeholder(slide, 0), title, size=28, bold=True)
    footer(slide, no)
    return slide


def add_navy_frame(slide):
    # Small editable accents that mimic the provided template.
    top = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(0), Inches(0.72), Inches(0.72))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()
    bot = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(12.61), Inches(6.78), Inches(0.72), Inches(0.72))
    bot.rotation = 180
    bot.fill.solid()
    bot.fill.fore_color.rgb = NAVY
    bot.line.fill.background()


def add_rect(slide, left, top, width, height, fill=WHITE, line=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_placeholder_text(placeholder(slide, 0), "CMTFlow：控制器引导的分层投资组合管理框架", size=30, color=BLACK, align=PP_ALIGN.CENTER)
    set_placeholder_text(
        placeholder(slide, 1),
        "Hierarchical Portfolio Management with Controller-Guided Base Revision and Daily Refinement\n作者 / 单位 / 日期",
        size=16,
        color=GRAY,
        bold=False,
        align=PP_ALIGN.CENTER,
    )
    add_navy_frame(slide)
    add_textbox(slide, "CMTFlow", 0.98, 1.14, 2.2, 0.36, size=18, color=RED, bold=True, font=FONT_EN)
    footer(slide, 1)


def build_background(prs):
    slide = title_slide(prs, 2, "研究背景：组合决策不是单一日频动作")
    add_navy_frame(slide)
    add_card(slide, "中期持仓", "形成稳定 base portfolio", 1.0, 2.05, 2.65, 1.15, SOFT_BLUE, BLUE, title_size=17, body_size=13)
    add_arrow(slide, 3.86, 2.42, 0.58, 0.3, MID_GRAY)
    add_card(slide, "每日修正", "围绕基准组合做局部微调", 4.65, 2.05, 2.65, 1.15, SOFT_GREEN, GREEN, title_size=17, body_size=13)
    add_arrow(slide, 7.52, 2.42, 0.58, 0.3, MID_GRAY)
    add_card(slide, "异常退出", "市场状态恶化时触发切换", 8.3, 2.05, 2.65, 1.15, SOFT_RED, RED, title_size=17, body_size=13)
    add_arrow(slide, 5.62, 3.55, 0.9, 0.36, MID_GRAY)
    add_textbox(slide, "核心挑战：稳定性与响应速度之间的权衡", 1.1, 4.75, 10.8, 0.45, size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def build_three_decisions(prs):
    slide = title_slide(prs, 3, "问题定义：本文关注的三个核心决策")
    add_navy_frame(slide)
    cards = [
        ("When", "当前 base portfolio 是否已经失效？\nController", RED, SOFT_RED),
        ("What", "如果切换，新的中期基准组合是什么？\nOuter Actor", BLUE, SOFT_BLUE),
        ("How", "在基准组合内部，如何每日权重微调？\nInner Actor", PURPLE, SOFT_PURPLE),
    ]
    for i, (title, body, color, fill) in enumerate(cards):
        add_card(slide, title, body, 0.95 + i * 4.05, 2.0, 3.1, 2.2, fill, color, title_size=24, body_size=14)
    add_textbox(slide, "不同时间尺度的决策，不应压缩进单一每日动作", 1.25, 5.35, 10.8, 0.42, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def build_method_flow(prs):
    slide = title_slide(prs, 4, "方法架构：从候选组合到最终执行权重")
    add_navy_frame(slide)
    items = [
        ("Market State", BLUE),
        ("Outer\n生成 w_t^{cand}", BLUE),
        ("Controller\nhold / switch", RED),
        ("Base Selector\n得到 b_t", BLUE),
        ("Inner\n输出 w_t", BLUE),
    ]
    for i, (text, color) in enumerate(items):
        add_card(slide, text, "", 0.65 + i * 2.52, 2.25, 1.82, 1.1, SOFT_RED if color == RED else SOFT_BLUE, color, title_size=14, body_size=1)
        if i < len(items) - 1:
            add_arrow(slide, 2.5 + i * 2.52, 2.62, 0.45, 0.28, MID_GRAY)
    add_card(slide, "关键", "Controller 比较旧 base 与候选 base，而不是机械按日历换仓", 1.0, 4.25, 4.75, 1.0, WHITE, RED, title_size=16, body_size=13)
    add_textbox(slide, "b_t = w_t^{cand} if switch, else \\tilde{b}_t    |    w_t = Inner(b_t, state_t)", 6.0, 4.45, 6.2, 0.42, size=15, color=NAVY, bold=True)


def build_results(prs):
    slide = title_slide(prs, 5, "数值结果：主实验性能对比")
    add_navy_frame(slide)
    add_rect(slide, 0.9, 1.55, 5.55, 4.3, WHITE, NAVY)
    add_textbox(slide, "Cumulative Wealth Curves", 1.15, 1.75, 4.9, 0.3, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 6.85, 1.55, 5.55, 4.3, WHITE, NAVY)
    add_textbox(slide, "Main Metrics: TR / Sharpe / MDD / CR", 7.1, 1.75, 4.9, 0.3, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # Schematic editable chart lines/bars.
    for y, color in [(4.6, RED), (4.9, BLUE), (5.15, GRAY)]:
        line = slide.shapes.add_connector(1, Inches(1.45), Inches(y), Inches(5.7), Inches(y - 1.15))
        line.line.color.rgb = color
        line.line.width = Inches(0.02)
    for i, color in enumerate([GRAY, RED, BLUE, GRAY, RED, BLUE, GRAY, RED]):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.35 + i * 0.55), Inches(4.85 - (i % 3) * 0.33), Inches(0.28), Inches(0.82 + (i % 3) * 0.33))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
    add_textbox(slide, "结论：CMTFlow 改善风险收益权衡，而不是追求所有单项指标绝对最优", 1.0, 6.15, 11.3, 0.38, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def build_conclusion(prs):
    slide = title_slide(prs, 6, "讨论与总结")
    add_navy_frame(slide)
    for i, (title, body, color, fill) in enumerate(
        [
            ("1. 问题重构", "base revision\nbase construction\ndaily refinement", RED, SOFT_RED),
            ("2. 核心机制", "Controller 将固定周期换仓\n变成可学习事件策略", BLUE, SOFT_BLUE),
            ("3. 实验结论", "更稳健的风险收益权衡\nController 是主要自适应来源", PURPLE, SOFT_PURPLE),
        ]
    ):
        add_card(slide, title, body, 0.95 + i * 4.05, 2.0, 3.1, 2.1, fill, color, title_size=17, body_size=13)
    add_card(slide, "未来工作", "更多市场与资产类别｜更复杂交易约束｜更稳健风险控制", 1.0, 5.2, 11.3, 0.78, SOFT_GREEN, GREEN, title_size=15, body_size=13)


def build():
    prs = Presentation(TEMPLATE)
    delete_all_slides(prs)
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    build_cover(prs)
    build_background(prs)
    build_three_decisions(prs)
    build_method_flow(prs)
    build_results(prs)
    build_conclusion(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
