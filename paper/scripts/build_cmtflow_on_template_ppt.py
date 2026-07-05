from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_cmtflow_editable_ppt import (
    BLACK,
    BLUE,
    FONT,
    GRAY,
    GREEN,
    LIGHT_GRAY,
    MID_GRAY,
    NAVY,
    PAPER,
    PURPLE,
    RED,
    SOFT_BLUE,
    SOFT_GREEN,
    SOFT_PURPLE,
    SOFT_RED,
    TEMPLATE,
    WHITE,
    add_arrow,
    add_bullets,
    add_card,
    add_notes,
    add_picture_fit,
    add_textbox,
    delete_all_slides,
    parse_speech,
    set_cell_text,
    set_run_font,
)


ROOT = Path(__file__).resolve().parents[2]
SPEECH = PAPER / "cmtflow_20min_speech.md"
OUT = PAPER / "CMTFlow_20min_on_template_draft.pptx"


def placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(idx)


def clear_placeholder(slide, idx):
    ph = placeholder(slide, idx)
    elm = ph._element
    elm.getparent().remove(elm)
    return ph.left, ph.top, ph.width, ph.height


def set_placeholder_text(ph, text, size=20, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        set_run_font(run, size=size, color=color, bold=bold)


def set_placeholder_bullets(ph, bullets, size=18, color=BLACK):
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        for run in p.runs:
            set_run_font(run, size=size, color=color)


def fill_footer(slide, n):
    for idx, text, align in [
        (10, "", PP_ALIGN.LEFT),
        (11, "CMTFlow", PP_ALIGN.CENTER),
        (12, f"{n:02d}", PP_ALIGN.RIGHT),
    ]:
        try:
            ph = placeholder(slide, idx)
        except KeyError:
            continue
        set_placeholder_text(ph, text, size=9, color=GRAY, align=align)


def slide_with_layout(prs, layout_idx, slide_no, title=None, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if title is not None:
        set_placeholder_text(placeholder(slide, 0), title, size=28 if layout_idx != 2 else 34, color=BLACK, bold=True, align=PP_ALIGN.CENTER if layout_idx in (0, 2) else PP_ALIGN.LEFT)
    fill_footer(slide, slide_no)
    if notes:
        add_notes(slide, notes)
    return slide


def add_picture_in_placeholder(slide, idx, path: Path, border=True):
    ph = placeholder(slide, idx)
    left, top, width, height = ph.left, ph.top, ph.width, ph.height
    elm = ph._element
    elm.getparent().remove(elm)
    return add_picture_fit(slide, path, left / 914400, top / 914400, width / 914400, height / 914400, border=border)


def build():
    notes = parse_speech(SPEECH)
    prs = Presentation(TEMPLATE)
    delete_all_slides(prs)
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # 1. Title slide: use template title layout.
    s = slide_with_layout(prs, 0, 1, "CMTFlow：控制器引导的分层投资组合管理框架", notes=notes.get(1, ""))
    set_placeholder_text(
        placeholder(s, 1),
        "Hierarchical Portfolio Management with Controller-Guided Base Revision and Daily Refinement\n作者 / 单位 / 日期",
        size=18,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    # 2-4. Background and problem pages: use title-content layout.
    s = slide_with_layout(prs, 1, 2, "研究背景：投资组合管理的基本目标", notes=notes.get(2, ""))
    set_placeholder_bullets(
        placeholder(s, 1),
        [
            "投资组合管理是在不确定市场中进行多资产资本配置",
            "目标不仅是提高收益，还要控制波动、回撤和交易成本",
            "强化学习适合长期序列决策，但金融市场状态持续变化",
        ],
        size=21,
    )

    s = slide_with_layout(prs, 1, 3, "研究背景：组合决策不是单一日频动作", notes=notes.get(3, ""))
    set_placeholder_bullets(
        placeholder(s, 1),
        [
            "实际投资包含多时间尺度：中期持仓、每日修正、异常状态退出",
            "固定周期再平衡稳定但反应慢",
            "纯日频调仓灵活但容易噪声化、换手过高",
            "因此需要把“持仓段”和“每日微调”分开建模",
        ],
        size=20,
    )

    s = slide_with_layout(prs, 1, 4, "问题定义：带漂移和交易成本的动态组合", notes=notes.get(4, ""))
    set_placeholder_bullets(
        placeholder(s, 1),
        [
            "昨日组合经过价格变化后会产生 drifted portfolio",
            "当日重新配置会带来换手和交易成本",
            "策略目标是在长期收益、风险和成本之间取得平衡",
            "评价指标包括 Total Return、Annual Return、Sharpe、MDD、CR",
        ],
        size=20,
    )

    # 5. Three decisions: use title-only layout and editable cards.
    s = slide_with_layout(prs, 5, 5, "问题定义：本文关注的三个核心决策", notes=notes.get(5, ""))
    for i, (title, body, color, fill) in enumerate(
        [
            ("When", "当前 base portfolio 是否已经失效？\n是否应该退出当前持仓段？", RED, SOFT_RED),
            ("What", "如果切换，下一段应该持有\n什么样的基准组合？", BLUE, SOFT_BLUE),
            ("How", "在选定基准组合内部，\n如何进行每日权重微调？", PURPLE, SOFT_PURPLE),
        ]
    ):
        add_card(s, title, body, 1.0 + i * 4.05, 2.1, 3.1, 2.05, fill, color, title_size=24, body_size=14)
        add_textbox(s, ["Controller", "Outer Actor", "Inner Actor"][i], 1.1 + i * 4.05, 4.38, 2.9, 0.35, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)

    # 6-7. Related work and challenge.
    s = slide_with_layout(prs, 3, 6, "相关工作：从静态优化到深度强化学习", notes=notes.get(6, ""))
    set_placeholder_bullets(
        placeholder(s, 1),
        ["传统优化：解释性强，但依赖静态假设或单期目标", "端到端 DRL：直接优化长期收益，但常压缩中期配置和每日执行"],
        size=18,
    )
    set_placeholder_bullets(
        placeholder(s, 2),
        ["层级/自适应方法：提升灵活性，但通常缺少显式 base portfolio 记忆", "本文切入点：学习基准组合什么时候应该被修正"],
        size=18,
    )

    s = slide_with_layout(prs, 3, 7, "挑战：直接应用 RL 仍存在非平凡限制", notes=notes.get(7, ""))
    set_placeholder_bullets(
        placeholder(s, 1),
        ["日频动作容易受到短期噪声驱动", "固定再平衡无法识别持仓状态是否已经恶化", "换仓需要同时比较旧组合、候选组合、成本和持仓年龄"],
        size=18,
    )
    set_placeholder_bullets(
        placeholder(s, 2),
        ["需要显式 Controller 学习 hold / switch", "将固定日历规则转化为状态依赖事件策略", "不是替代选股模块，而是判断当前 base 是否仍值得持有"],
        size=18,
    )

    # 8. Overview: use template picture-title layout.
    s = slide_with_layout(prs, 8, 8, "方案总览：CMTFlow 的统一分层结构", notes=notes.get(8, ""))
    set_placeholder_bullets(
        placeholder(s, 2),
        ["Outer：生成候选 base portfolio", "Controller：判断 hold 或 switch", "Inner：围绕当前 base 做每日微调", "环境反馈收益、成本和训练信号"],
        size=16,
    )
    add_picture_in_placeholder(s, 1, PAPER / "figures/cmtflow_architecture_vector.png")

    # 9. Formulation.
    s = slide_with_layout(prs, 1, 9, "强化学习建模：状态、动作与奖励", notes=notes.get(9, ""))
    set_placeholder_bullets(
        placeholder(s, 1),
        [
            "State：近期市场张量、漂移当前持仓、候选组合、holding-state/action-comparison features",
            "Action：候选基准组合、hold/switch 决策、最终执行权重",
            "Reward：组合 log return 扣除交易成本",
            "目标：优化完整投资路径上的风险收益质量",
        ],
        size=20,
    )

    # 10-14. Method and training pages.
    for no, title, img, bullets in [
        (
            10,
            "方法架构：从候选组合到最终执行权重",
            PAPER / "figures/cmtflow_decision_flow_imagegen.png",
            ["Outer 生成候选组合 w_t^{cand}", "Controller 输出 switch probability", "Base selector 得到实际 b_t", "Inner 输出最终执行权重 w_t"],
        ),
        (
            11,
            "方法细节：Outer Actor 生成中期基准组合",
            PAPER / "figures/cmtflow_architecture_vector.png",
            ["面向持仓段级别决策，而不是每日最终交易", "输入长期市场窗口和当前持仓漂移状态", "输出稀疏 top-K candidate base portfolio", "回答 What to hold next"],
        ),
        (
            12,
            "方法细节：Controller 学习何时替换基准组合",
            PAPER / "figures/cmtflow_decision_flow_imagegen.png",
            ["输入近期市场张量、漂移当前持仓、候选组合和比较特征", "输出 exit probability 或 hold/switch 决策", "每日检查，30 天最大持仓上限", "把固定日历调仓变成事件策略"],
        ),
        (
            13,
            "方法细节：Inner Actor 做每日局部微调",
            PAPER / "figures/cmtflow_architecture_vector.png",
            ["不重新决定股票池，而是在 active base 内调整权重", "提供日频灵活性，但不脱离中期基准组合", "作用是局部 refinement，而不是主要换仓控制器"],
        ),
        (
            14,
            "模型训练：固定 HRL 预训练与 Controller 学习",
            PAPER / "figures/cmtflow_training_flow_vector.png",
            ["先训练 fixed-segment HRL backbone", "再训练每日 hold/switch Controller", "最终测试不使用固定再平衡周期", "PPO 与辅助信号共同优化策略"],
        ),
    ]:
        s = slide_with_layout(prs, 8, no, title, notes=notes.get(no, ""))
        set_placeholder_bullets(placeholder(s, 2), bullets, size=15 if no in (10, 12, 14) else 16)
        add_picture_in_placeholder(s, 1, img)

    # 15. Experimental setup.
    s = slide_with_layout(prs, 5, 15, "实验设置", notes=notes.get(15, ""))
    table = s.shapes.add_table(3, 5, Inches(0.92), Inches(1.9), Inches(7.75), Inches(1.55)).table
    headers = ["Market", "#Stocks", "Train", "Valid", "Test"]
    rows = [
        ["Nasdaq-100", "39", "2000-04-07\n2017-12-29", "2018-01-02\n2020-04-22", "2020-04-23\n2025-10-03"],
        ["CSI-300", "53", "2000-04-07\n2017-12-28", "2018-01-02\n2019-12-31", "2020-01-02\n2025-02-28"],
    ]
    for j, h in enumerate(headers):
        set_cell_text(table.cell(0, j), h, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        table.cell(0, j).fill.solid()
        table.cell(0, j).fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            set_cell_text(table.cell(i, j), val, size=9.5, align=PP_ALIGN.CENTER)
            table.cell(i, j).fill.solid()
            table.cell(i, j).fill.fore_color.rgb = WHITE if i == 1 else LIGHT_GRAY
    add_card(s, "Baselines", "传统策略\n深度 RL baseline\n固定 5/10/20/30/60 天 controller", 9.05, 1.85, 3.15, 1.55, SOFT_BLUE, BLUE)
    add_card(s, "Metrics", "TR / AR / Vol\nSharpe ratio\nMDD / CR", 9.05, 3.7, 3.15, 1.55, SOFT_PURPLE, PURPLE)

    # 16. Main result: use title-only template layout for two large figures.
    s = slide_with_layout(prs, 5, 16, "数值结果：主实验性能对比", notes=notes.get(16, ""))
    add_picture_fit(s, PAPER / "figures/main_equity_curves.png", 0.88, 1.55, 5.75, 4.75)
    add_picture_fit(s, PAPER / "figures/main_metric_bars.png", 6.85, 1.55, 5.55, 4.75)
    add_textbox(s, "Nasdaq: TR 265.53%, MDD 18.62%  |  CSI-300: TR 204.99%, Sharpe 1.14, MDD 22.78%", 0.92, 6.35, 11.5, 0.35, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 17. Ablation: picture-title template.
    s = slide_with_layout(prs, 8, 17, "数值结果：消融实验与机制验证", notes=notes.get(17, ""))
    set_placeholder_bullets(
        placeholder(s, 2),
        ["Learned controller 优于固定窗口切换策略", "Controller 是主要自适应来源", "Inner Actor 更偏局部调权和风险控制", "提升来自学习何时修正 base portfolio"],
        size=15,
    )
    add_picture_in_placeholder(s, 1, PAPER / "figures/ablation_metric_bars.png")

    # 18. Case study: title-only template with three images.
    s = slide_with_layout(prs, 5, 18, "案例研究：Controller 与 Inner 的可解释行为", notes=notes.get(18, ""))
    add_picture_fit(s, PAPER / "figures/explainability/controller_switch_cases.png", 0.85, 1.5, 3.9, 4.75)
    add_picture_fit(s, PAPER / "figures/explainability/random_switch_comparison.png", 4.9, 1.5, 3.75, 4.75)
    add_picture_fit(s, PAPER / "figures/explainability/inner_actor_base_adjustment.png", 8.85, 1.5, 3.65, 4.75)
    add_textbox(s, "Controller 解释换仓时机；随机切换说明不是“切换本身”带来收益；Inner 展示局部 tilt。", 0.95, 6.35, 11.35, 0.35, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 19. Conclusion.
    s = slide_with_layout(prs, 1, 19, "讨论与总结", notes=notes.get(19, ""))
    set_placeholder_bullets(
        placeholder(s, 1),
        [
            "CMTFlow 将组合管理拆成 base revision、base construction 和 daily refinement",
            "Controller 让换仓从固定周期规则变成可学习事件策略",
            "实验和消融表明 controller 是动态修正行为的主要来源",
            "未来工作：更强市场泛化、更复杂交易约束、更稳健风险控制",
        ],
        size=20,
    )

    # 20. Q&A: use template section-title layout.
    s = slide_with_layout(prs, 2, 20, "谢谢，欢迎提问", notes=notes.get(20, ""))
    set_placeholder_text(
        placeholder(s, 1),
        "Learning when to revise is as important as learning what to hold.",
        size=18,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
